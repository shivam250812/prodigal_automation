import os
import json
import logging
from typing import List, Dict, Any
from pathlib import Path
from dotenv import load_dotenv
import openai
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain.prompts import ChatPromptTemplate
import pytesseract
import easyocr
from pdf2image import convert_from_path
import cv2
import numpy as np
from pptx import Presentation
import tempfile

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Set OpenAI API key
openai.api_key = os.getenv("OPENAI_API_KEY")

# Set Tesseract path for macOS
pytesseract.pytesseract.tesseract_cmd = '/opt/homebrew/bin/tesseract'

class MultiFileRAGSystem:
    """RAG system that processes multiple files and allows post-generation manual review."""
    
    def __init__(self):
        """Initialize the RAG system."""
        self.embeddings = OpenAIEmbeddings()
        self.llm = ChatOpenAI(model="gpt-4o", temperature=0.1)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        
    def get_loader(self, file_path: str):
        """Get appropriate loader based on file extension."""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == '.txt':
            from langchain_community.document_loaders import TextLoader
            return TextLoader(file_path)
        elif file_ext == '.pdf':
            from langchain_community.document_loaders import PyPDFLoader
            return PyPDFLoader(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self._create_pptx_loader(file_path)
        elif file_ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            return self._create_image_loader(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    
    def _create_pptx_loader(self, file_path: str):
        """Create custom PPTX loader."""
        class PPTXLoader:
            def __init__(self, file_path):
                self.file_path = file_path
            
            def load(self):
                """Load PPTX file and extract text from slides and tables."""
                try:
                    prs = Presentation(self.file_path)
                    documents = []
                    
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_text = f"Slide {slide_num}:\n"
                        
                        # Extract text from shapes
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text += shape.text + "\n"
                        
                        # Extract text from tables
                        for shape in slide.shapes:
                            if shape.has_table:
                                table = shape.table
                                for row in table.rows:
                                    row_text = " | ".join([cell.text for cell in row.cells])
                                    slide_text += row_text + "\n"
                        
                        if slide_text.strip() != f"Slide {slide_num}:\n":
                            documents.append(Document(
                                page_content=slide_text.strip(),
                                metadata={"source": self.file_path, "slide": slide_num}
                            ))
                    
                    return documents
                except Exception as e:
                    logger.error(f"Error processing PPTX file {self.file_path}: {e}")
                    return []
        
        return PPTXLoader(file_path)
    
    def _create_image_loader(self, file_path: str):
        """Create custom image loader with OCR."""
        class ImageLoader:
            def __init__(self, file_path):
                self.file_path = file_path
            
            def load(self):
                """Load image and perform OCR."""
                try:
                    text = self._perform_ocr(self.file_path)
                    if text.strip():
                        return [Document(
                            page_content=text,
                            metadata={"source": self.file_path, "type": "image"}
                        )]
                    return []
                except Exception as e:
                    logger.error(f"Error processing image {self.file_path}: {e}")
                    return []
            
            def _perform_ocr(self, image_path: str) -> str:
                """Perform OCR on image using both EasyOCR and Tesseract."""
                try:
                    # Try EasyOCR first
                    reader = easyocr.Reader(['en'])
                    results = reader.readtext(image_path)
                    easyocr_text = ' '.join([text[1] for text in results])
                    
                    # Try Tesseract
                    tesseract_text = pytesseract.image_to_string(image_path)
                    
                    # Return the longer result (usually more accurate)
                    if len(easyocr_text) > len(tesseract_text):
                        return easyocr_text
                    else:
                        return tesseract_text
                        
                except Exception as e:
                    logger.error(f"OCR error for {image_path}: {e}")
                    return ""
        
        return ImageLoader(file_path)
    
    def process_pdf_with_images(self, pdf_path: str):
        """Process PDF and extract text from embedded images."""
        try:
            # Convert PDF pages to images
            images = convert_from_path(pdf_path)
            documents = []
            
            for page_num, image in enumerate(images, 1):
                # Save image temporarily
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                    image.save(tmp_file.name, 'PNG')
                    tmp_path = tmp_file.name
                
                try:
                    # Perform OCR on the image
                    text = self._perform_ocr(tmp_path)
                    if text.strip():
                        documents.append(Document(
                            page_content=f"Page {page_num} (Image): {text}",
                            metadata={"source": pdf_path, "page": page_num, "type": "pdf_image"}
                        ))
                finally:
                    # Clean up temporary file
                    os.unlink(tmp_path)
            
            return documents
        except Exception as e:
            logger.error(f"Error processing PDF images {pdf_path}: {e}")
            return []
    
    def _perform_ocr(self, image_path: str) -> str:
        """Perform OCR on image."""
        try:
            # Try EasyOCR first
            reader = easyocr.Reader(['en'])
            results = reader.readtext(image_path)
            easyocr_text = ' '.join([text[1] for text in results])
            
            # Try Tesseract
            tesseract_text = pytesseract.image_to_string(image_path)
            
            # Return the longer result
            if len(easyocr_text) > len(tesseract_text):
                return easyocr_text
            else:
                return tesseract_text
                
        except Exception as e:
            logger.error(f"OCR error for {image_path}: {e}")
            return ""
    
    def process_documents(self, file_paths: List[str]) -> List[Any]:
        """Process multiple documents and return chunks."""
        all_chunks = []
        processed_files = []
        
        for file_path in file_paths:
            try:
                print(f"📄 Processing: {Path(file_path).name}")
                
                # Get appropriate loader
                loader = self.get_loader(file_path)
                
                # Load documents
                if hasattr(loader, 'load'):
                    documents = loader.load()
                else:
                    documents = loader.load()
                
                # Process PDF images if it's a PDF
                if Path(file_path).suffix.lower() == '.pdf':
                    print(f"  🔍 Processing PDF for embedded images...")
                    image_documents = self.process_pdf_with_images(file_path)
                    if image_documents:
                        documents.extend(image_documents)
                        print(f"  ✅ Added {len(image_documents)} image-based documents")
                
                # Split documents into chunks
                for doc in documents:
                    chunks = self.text_splitter.split_documents([doc])
                    all_chunks.extend(chunks)
                
                processed_files.append(file_path)
                print(f"  ✅ Processed {len(documents)} documents")
                
            except Exception as e:
                logger.error(f"Error processing {file_path}: {e}")
                print(f"  ❌ Failed to process {Path(file_path).name}")
        
        print(f"✅ Total chunks: {len(all_chunks)}")
        return all_chunks, processed_files
    
    def create_knowledge_base(self, chunks: List[Any]):
        """Create knowledge base from chunks."""
        print("🧠 Creating knowledge base...")
        
        try:
            # Test embeddings
            test_text = "test"
            test_embedding = self.embeddings.embed_query(test_text)
            print(f"✅ Embedding test successful. Vector dimension: {len(test_embedding)}")
            
            # Create vector store
            vectorstore = FAISS.from_documents(chunks, self.embeddings)
            retriever = vectorstore.as_retriever(search_kwargs={"k": 5})
            
            print("✅ Vector store created successfully")
            return retriever
            
        except Exception as e:
            logger.error(f"Error creating knowledge base: {e}")
            raise
    
    def generate_all_answers(self, questions: List[str], retriever) -> List[Dict]:
        """Generate answers for all questions at once."""
        print(f"🤖 Generating answers for {len(questions)} questions...")
        
        qa_pairs = []
        
        for i, question in enumerate(questions, 1):
            try:
                print(f"  Generating answer {i}/{len(questions)}: {question[:50]}...")
                
                # Retrieve relevant documents
                docs = retriever.get_relevant_documents(question)
                context = "\n\n".join([doc.page_content for doc in docs])
                
                # Create prompt
                prompt = ChatPromptTemplate.from_template("""
                Based on the following context, answer the question accurately and comprehensively.
                
                Context:
                {context}
                
                Question: {question}
                
                Answer:""")
                
                # Generate answer
                chain = prompt | self.llm
                response = chain.invoke({"context": context, "question": question})
                answer = response.content.strip()
                
                # Create Q&A pair
                qa_pair = {
                    "question": question,
                    "ai_answer": answer,
                    "manual_answer": None,  # Will be filled during review
                    "has_manual_input": False,
                    "confidence_score": self._calculate_confidence(answer),
                    "edited": False
                }
                
                qa_pairs.append(qa_pair)
                print(f"    ✅ Generated successfully")
                
            except Exception as e:
                logger.error(f"Error generating answer for question {i}: {e}")
                print(f"    ❌ Failed to generate answer")
                
                # Add empty Q&A pair
                qa_pairs.append({
                    "question": question,
                    "ai_answer": "Error generating answer",
                    "manual_answer": None,
                    "has_manual_input": False,
                    "confidence_score": 0.0,
                    "edited": False
                })
        
        print(f"✅ Generated {len(qa_pairs)} answers")
        return qa_pairs
    
    def _calculate_confidence(self, answer: str) -> float:
        """Calculate confidence score for an answer."""
        # Simple confidence calculation based on answer length and content
        if not answer or answer.strip() == "":
            return 0.0
        
        # Check for low confidence indicators
        low_confidence_indicators = [
            "i don't know", "cannot determine", "not provided", "not mentioned",
            "no information", "unclear", "not specified", "not available"
        ]
        
        answer_lower = answer.lower()
        for indicator in low_confidence_indicators:
            if indicator in answer_lower:
                return 0.3
        
        # Base confidence on answer length and quality
        if len(answer.strip()) < 50:
            return 0.4
        elif len(answer.strip()) < 100:
            return 0.6
        else:
            return 0.8
    
    def review_answers(self, qa_pairs: List[Dict]) -> List[Dict]:
        """Review all answers after generation."""
        print("\n🚀 Starting Post-Generation Review")
        print(f"📊 Total Q&A pairs to review: {len(qa_pairs)}")
        print(f"💡 You can accept, edit, or skip each answer")
        print(f"💡 Type 'q' at any time to quit the review")
        
        reviewed_pairs = []
        skipped_count = 0
        edited_count = 0
        accepted_count = 0
        
        for i, qa_pair in enumerate(qa_pairs):
            self._display_qa_pair(qa_pair, i, len(qa_pairs))
            
            while True:
                choice = self._get_user_choice()
                
                if choice == 'a':  # Accept
                    reviewed_pairs.append(qa_pair)
                    accepted_count += 1
                    print(f"✅ Accepted AI answer")
                    break
                
                elif choice == 'e':  # Edit
                    edited_pair = self._edit_answer(qa_pair)
                    reviewed_pairs.append(edited_pair)
                    edited_count += 1
                    break
                
                elif choice == 's':  # Skip (keep AI answer)
                    reviewed_pairs.append(qa_pair)
                    skipped_count += 1
                    print(f"⏭️  Skipped (keeping AI answer)")
                    break
                
                elif choice == 'q':  # Quit
                    print(f"\n🛑 Review stopped early")
                    # Add remaining pairs as accepted
                    reviewed_pairs.extend(qa_pairs[i:])
                    return reviewed_pairs
                
                elif choice == 'h':  # Help
                    self._show_help()
                
                else:
                    print("❌ Invalid choice. Please try again.")
        
        print(f"\n🎉 Review Complete!")
        print(f"📊 Summary:")
        print(f"  ✅ Accepted: {accepted_count}")
        print(f"  ✏️  Edited: {edited_count}")
        print(f"  ⏭️  Skipped: {skipped_count}")
        print(f"  📋 Total: {len(reviewed_pairs)}")
        
        return reviewed_pairs
    
    def _display_qa_pair(self, qa_pair: Dict, index: int, total: int):
        """Display a Q&A pair for review."""
        print(f"\n{'='*80}")
        print(f"📋 Q&A Pair {index + 1} of {total}")
        print(f"{'='*80}")
        
        print(f"❓ Question: {qa_pair.get('question', 'N/A')}")
        print(f"\n🤖 AI Answer: {qa_pair.get('ai_answer', 'N/A')}")
        
        if qa_pair.get('confidence_score'):
            print(f"\n📊 Confidence Score: {qa_pair['confidence_score']:.2f}")
        
        if qa_pair.get('has_manual_input'):
            print(f"\n✏️  Manual Answer: {qa_pair['manual_answer']}")
    
    def _get_user_choice(self) -> str:
        """Get user choice for Q&A pair review."""
        print(f"\n📝 Review Options:")
        print("  [a] Accept AI answer")
        print("  [e] Edit answer manually")
        print("  [s] Skip (keep AI answer)")
        print("  [q] Quit review")
        print("  [h] Help")
        
        choice = input("\nEnter your choice: ").strip().lower()
        return choice
    
    def _edit_answer(self, qa_pair: Dict) -> Dict:
        """Edit the answer for a Q&A pair."""
        print(f"\n✏️  Editing answer for: {qa_pair['question'][:100]}...")
        
        current_answer = qa_pair.get('ai_answer', '')
        print(f"\nCurrent AI answer: {current_answer}")
        
        print(f"\nEnter your manual answer:")
        manual_answer = input("Manual answer: ").strip()
        
        if manual_answer:
            qa_pair['manual_answer'] = manual_answer
            qa_pair['has_manual_input'] = True
            qa_pair['edited'] = True
            print(f"✅ Manual answer saved")
        else:
            print(f"⏭️  Keeping AI answer")
        
        return qa_pair
    
    def _show_help(self):
        """Show help information."""
        print(f"\n📖 Help:")
        print("  [a] Accept - Keep the AI-generated answer")
        print("  [e] Edit - Enter your own manual answer")
        print("  [s] Skip - Keep the AI answer (same as accept)")
        print("  [q] Quit - Stop reviewing and save current progress")
        print("  [h] Help - Show this help message")
    
    def save_qa_pairs(self, qa_pairs: List[Dict], filename: str = "qa_pairs.json"):
        """Save Q&A pairs to JSON file."""
        try:
            with open(filename, 'w', encoding='utf-8') as f:
                json.dump(qa_pairs, f, indent=2, ensure_ascii=False)
            print(f"✅ Saved {len(qa_pairs)} Q&A pairs to {filename}")
        except Exception as e:
            logger.error(f"Error saving Q&A pairs: {e}")
            print(f"❌ Error saving Q&A pairs")
    
    def load_qa_pairs(self, filename: str = "qa_pairs.json") -> List[Dict]:
        """Load Q&A pairs from JSON file."""
        try:
            if os.path.exists(filename):
                with open(filename, 'r', encoding='utf-8') as f:
                    return json.load(f)
            return []
        except Exception as e:
            logger.error(f"Error loading Q&A pairs: {e}")
            return []
    
    def load_questions(self, filename: str = "questions.json") -> List[str]:
        """Load questions from JSON file."""
        try:
            if os.path.exists(filename):
                with open(filename, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    if isinstance(data, list):
                        return [item.get('question', item) if isinstance(item, dict) else str(item) for item in data]
                    else:
                        return []
            else:
                print(f"❌ Questions file not found: {filename}")
                return []
        except Exception as e:
            logger.error(f"Error loading questions: {e}")
            return []

def main():
    """Main function for the multi-file RAG system."""
    print("🚀 Multi-File RAG System")
    print("=" * 50)
    print("Features: Multiple files → Generate all answers → Manual review")
    print("=" * 50)
    
    rag_system = MultiFileRAGSystem()
    
    print("\n📋 Choose operation:")
    print("1. Process multiple files and generate answers")
    print("2. Review existing Q&A pairs")
    print("3. Exit")
    
    choice = input("\nEnter your choice (1-3): ").strip()
    
    if choice == "1":
        # Get file paths
        print("\n📁 Enter file paths (one per line, press Enter twice when done):")
        file_paths = []
        while True:
            file_path = input("File path: ").strip()
            if not file_path:
                break
            if os.path.isfile(file_path):
                file_paths.append(file_path)
            else:
                print(f"❌ File not found: {file_path}")
        
        if not file_paths:
            print("❌ No valid files provided.")
            return
        
        # Process documents
        chunks, processed_files = rag_system.process_documents(file_paths)
        
        if not chunks:
            print("❌ No documents were processed successfully.")
            return
        
        # Create knowledge base
        retriever = rag_system.create_knowledge_base(chunks)
        
        # Load questions
        questions = rag_system.load_questions()
        if not questions:
            print("❌ No questions found. Please create a questions.json file.")
            return
        
        # Generate all answers first
        qa_pairs = rag_system.generate_all_answers(questions, retriever)
        
        # Save generated pairs
        rag_system.save_qa_pairs(qa_pairs, "generated_qa_pairs.json")
        
        # Ask if user wants to review
        review_choice = input("\n🤔 Would you like to review the answers now? (y/n): ").strip().lower()
        if review_choice in ['y', 'yes']:
            reviewed_pairs = rag_system.review_answers(qa_pairs)
            rag_system.save_qa_pairs(reviewed_pairs, "reviewed_qa_pairs.json")
        else:
            print("✅ Answers saved without review")
    
    elif choice == "2":
        # Review existing Q&A pairs
        qa_pairs = rag_system.load_qa_pairs()
        
        if not qa_pairs:
            print("❌ No Q&A pairs found to review.")
            return
        
        reviewed_pairs = rag_system.review_answers(qa_pairs)
        rag_system.save_qa_pairs(reviewed_pairs, "reviewed_qa_pairs.json")
    
    elif choice == "3":
        print("👋 Goodbye!")
    
    else:
        print("❌ Invalid choice.")

if __name__ == "__main__":
    main() 